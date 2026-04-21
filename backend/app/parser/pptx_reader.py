from __future__ import annotations

import json
import platform
import re
import shutil
import subprocess
import tempfile
from collections.abc import Callable
from io import BytesIO
from pathlib import Path

import fitz
from pptx import Presentation

from backend.app.common.exceptions import ApiError
from backend.app.parser.schemas import ExtractedPresentation, ExtractedSlide, FileInfo
from backend.app.parser.source_loader import load_source_bytes

POWERSHELL_EXE = Path(r"C:\Windows\System32\WindowsPowerShell\v1.0\powershell.exe")
LIBREOFFICE_CANDIDATES = (
    Path("/Applications/LibreOffice.app/Contents/MacOS/soffice"),
    Path("/usr/bin/soffice"),
    Path("/usr/local/bin/soffice"),
)
MACOS_POWERPOINT_BUNDLE_ID = "com.microsoft.Powerpoint"
MACOS_POWERPOINT_APP_CANDIDATES = (
    Path("/Applications/Microsoft PowerPoint.app"),
)
OSASCRIPT_EXE = Path("/usr/bin/osascript")
POWERPOINT_EXPORT_SCRIPT = """\
param(
    [Parameter(Mandatory = $true)][string]$InputPath,
    [Parameter(Mandatory = $true)][string]$OutputDir
)
$ErrorActionPreference = 'Stop'
$ppt = $null
$presentation = $null
try {
    $ppt = New-Object -ComObject PowerPoint.Application
    $presentation = $ppt.Presentations.Open($InputPath, $true, $false, $false)
    $presentation.Export($OutputDir, 'PNG')
}
finally {
    if ($presentation -ne $null) {
        $presentation.Close()
        [void][System.Runtime.InteropServices.Marshal]::ReleaseComObject($presentation)
    }
    if ($ppt -ne $null) {
        $ppt.Quit()
        [void][System.Runtime.InteropServices.Marshal]::ReleaseComObject($ppt)
    }
}
"""


def extract_pptx_presentation(
    file_url: str,
    *,
    preview_output_dir: Path | None = None,
    preview_public_base: str | None = None,
) -> tuple[FileInfo, ExtractedPresentation]:
    file_name, file_bytes = load_source_bytes(file_url, file_label="PPTX", default_name="courseware.pptx")
    normalized_name = file_name.lower()
    if normalized_name.endswith(".ppt"):
        raise ApiError(code=400, msg="当前环境暂仅支持 .pptx 文件预览生成，请先将 .ppt 转为 .pptx 后重试", status_code=400)
    if not normalized_name.endswith(".pptx"):
        raise ApiError(code=400, msg="当前 demo 仅支持 .pptx 文件，请传入可访问的 PPTX 文件", status_code=400)

    try:
        presentation = Presentation(BytesIO(file_bytes))
    except Exception as exc:  # noqa: BLE001
        raise ApiError(code=400, msg="PPTX 文件无法解析，请确认文件内容有效", status_code=400) from exc

    slide_count = len(presentation.slides)
    preview_urls = _build_preview_urls(
        file_name=file_name,
        file_bytes=file_bytes,
        slide_count=slide_count,
        output_dir=preview_output_dir,
        public_base=preview_public_base,
    )

    slides: list[ExtractedSlide] = []
    for index, slide in enumerate(presentation.slides, start=1):
        body_texts: list[str] = []
        table_texts: list[str] = []
        title: str | None = None

        for shape in slide.shapes:
            current_title = _extract_title(shape)
            if current_title and not title:
                title = current_title

            text_frame = getattr(shape, "text_frame", None)
            if text_frame is not None:
                text = _normalize_text(getattr(text_frame, "text", None))
                if text:
                    body_texts.append(text)

            table = getattr(shape, "table", None)
            if table is not None:
                rows: list[str] = []
                for row in table.rows:
                    values = [_normalize_text(cell.text) for cell in row.cells]
                    values = [value for value in values if value]
                    if values:
                        rows.append(" | ".join(values))
                if rows:
                    table_texts.append("\n".join(rows))

        notes = _extract_notes(slide)
        if not title:
            title = body_texts[0][:40] if body_texts else f"第{index}页"

        slides.append(
            ExtractedSlide(
                slideNumber=index,
                title=title,
                bodyTexts=body_texts,
                tableTexts=table_texts,
                notes=notes,
                previewUrl=preview_urls.get(index),
            )
        )

    if not slides:
        raise ApiError(code=400, msg="PPTX 中没有可解析的页面", status_code=400)

    file_info = FileInfo(fileName=file_name, fileSize=len(file_bytes), pageCount=len(slides))
    extracted = ExtractedPresentation(sourceType="pptx", slides=slides)
    return file_info, extracted


def _build_preview_urls(
    *,
    file_name: str,
    file_bytes: bytes,
    slide_count: int,
    output_dir: Path | None,
    public_base: str | None,
) -> dict[int, str]:
    if output_dir is None or not public_base:
        return {}

    output_dir.mkdir(parents=True, exist_ok=True)
    with tempfile.TemporaryDirectory(prefix="pptx-preview-", dir="/tmp") as temp_dir_name:
        temp_dir = Path(temp_dir_name)
        input_path = temp_dir / Path(file_name).name
        export_dir = temp_dir / "exported"
        export_dir.mkdir(parents=True, exist_ok=True)
        input_path.write_bytes(file_bytes)

        exported_files: list[Path] = []
        export_errors: list[str] = []
        for exporter_name, exporter in _resolve_preview_exporters():
            _clear_exported_images(export_dir)
            try:
                exporter(input_path, export_dir, slide_count)
                exported_files = _collect_exported_images(export_dir)
                if not exported_files:
                    raise ApiError(code=500, msg=f"{exporter_name} 未导出任何页图，请检查本机渲染能力", status_code=500)
                if slide_count and len(exported_files) != slide_count:
                    raise ApiError(
                        code=500,
                        msg=f"{exporter_name} 导出的页图数量异常，期望 {slide_count} 页，实际 {len(exported_files)} 页",
                        status_code=500,
                    )
                break
            except ApiError as exc:
                export_errors.append(exc.msg)
                continue

        if not exported_files:
            message = export_errors[-1] if export_errors else "当前环境没有可用的课件页图导出器"
            raise ApiError(code=500, msg=message, status_code=500)

        preview_urls: dict[int, str] = {}
        for index, source_path in enumerate(exported_files, start=1):
            target_path = output_dir / f"page-{index}.png"
            shutil.copy2(source_path, target_path)
            preview_urls[index] = f"{public_base.rstrip('/')}/page-{index}.png"
        return preview_urls


def _resolve_preview_exporters() -> list[tuple[str, Callable[[Path, Path, int], None]]]:
    exporters: list[tuple[str, Callable[[Path, Path, int], None]]] = []
    current_system = platform.system().lower()
    if current_system == "windows" and POWERSHELL_EXE.exists():
        exporters.append(("PowerPoint", _export_with_powerpoint))
    elif current_system == "darwin":
        macos_powerpoint_app = _find_macos_powerpoint_app()
        if macos_powerpoint_app is not None and OSASCRIPT_EXE.exists():
            exporters.append(
                (
                    "PowerPoint for Mac",
                    lambda input_path, output_dir, slide_count: _export_with_macos_powerpoint(
                        macos_powerpoint_app,
                        input_path,
                        output_dir,
                        slide_count,
                    ),
                )
            )

    soffice_path = _find_soffice_executable()
    if soffice_path is not None:
        exporters.append(
            (
                "LibreOffice",
                lambda input_path, output_dir, slide_count: _export_with_libreoffice(
                    soffice_path,
                    input_path,
                    output_dir,
                    slide_count,
                ),
            )
        )

    if exporters:
        return exporters

    if current_system == "windows":
        raise ApiError(
            code=500,
            msg="当前环境既缺少 PowerPoint/PowerShell，也未找到 LibreOffice soffice，无法导出课件页图",
            status_code=500,
        )
    if current_system == "darwin":
        raise ApiError(
            code=500,
            msg="当前环境既未找到 PowerPoint for Mac，也未找到 LibreOffice soffice，无法导出课件页图",
            status_code=500,
        )
    raise ApiError(
        code=500,
        msg="当前环境未找到 LibreOffice soffice，macOS/Linux 请先安装 LibreOffice 后再导出课件页图",
        status_code=500,
    )


def _find_macos_powerpoint_app() -> Path | None:
    for candidate in MACOS_POWERPOINT_APP_CANDIDATES:
        if candidate.exists():
            return candidate
    try:
        completed = subprocess.run(
            ["mdfind", f"kMDItemCFBundleIdentifier == '{MACOS_POWERPOINT_BUNDLE_ID}'"],
            capture_output=True,
            text=True,
            check=False,
        )
    except OSError:
        return None
    if completed.returncode != 0:
        return None
    for line in completed.stdout.splitlines():
        candidate = Path(line.strip())
        if candidate.exists():
            return candidate
    return None


def _find_soffice_executable() -> Path | None:
    soffice_path = shutil.which("soffice")
    if soffice_path:
        return Path(soffice_path)
    for candidate in LIBREOFFICE_CANDIDATES:
        if candidate.exists():
            return candidate
    return None


def _export_with_powerpoint(input_path: Path, output_dir: Path, _slide_count: int) -> None:
    if not POWERSHELL_EXE.exists():
        raise ApiError(code=500, msg="当前环境缺少 PowerShell，无法调用 PowerPoint 导出课件页图", status_code=500)

    script_path = output_dir.parent / "export_pptx_preview.ps1"
    script_path.write_text(POWERPOINT_EXPORT_SCRIPT, encoding="utf-8")
    completed = subprocess.run(
        [
            str(POWERSHELL_EXE),
            "-NoProfile",
            "-NonInteractive",
            "-ExecutionPolicy",
            "Bypass",
            "-File",
            str(script_path),
            "-InputPath",
            str(input_path),
            "-OutputDir",
            str(output_dir),
        ],
        capture_output=True,
        text=True,
        check=False,
    )
    if completed.returncode != 0:
        detail = (completed.stderr or completed.stdout or "").strip().splitlines()
        message = detail[-1] if detail else "未知错误"
        raise ApiError(code=500, msg=f"调用 PowerPoint 导出页图失败：{message}", status_code=500)


def _export_with_macos_powerpoint(_powerpoint_app: Path, input_path: Path, output_dir: Path, slide_count: int) -> None:
    if not OSASCRIPT_EXE.exists():
        raise ApiError(code=500, msg="当前环境缺少 osascript，无法调用 PowerPoint for Mac 导出课件页图", status_code=500)

    pdf_path = output_dir.parent / "exported-preview.pdf"
    if pdf_path.exists():
        pdf_path.unlink()
    input_literal = _to_applescript_posix_file_literal(input_path)
    output_literal = _to_applescript_posix_file_literal(pdf_path)
    completed = subprocess.run(
        [
            str(OSASCRIPT_EXE),
            "-e",
            'tell application "Microsoft PowerPoint" to activate',
            "-e",
            f'tell application "Microsoft PowerPoint" to open {input_literal}',
            "-e",
            f'tell application "Microsoft PowerPoint" to save active presentation in {output_literal} as save as PDF',
            "-e",
            'tell application "Microsoft PowerPoint" to close active presentation saving no',
        ],
        capture_output=True,
        text=True,
        check=False,
    )
    if completed.returncode != 0:
        detail = (completed.stderr or completed.stdout or "").strip().splitlines()
        message = detail[-1] if detail else "未知错误"
        raise ApiError(code=500, msg=f"调用 PowerPoint for Mac 导出 PDF 失败：{message}", status_code=500)
    if not pdf_path.exists():
        raise ApiError(code=500, msg="PowerPoint for Mac 未导出 PDF 文件，请检查本机 Office 渲染能力", status_code=500)

    _render_pdf_to_png(pdf_path, output_dir, slide_count=slide_count)


def _to_applescript_posix_file_literal(path: Path) -> str:
    escaped = str(path).replace("\\", "\\\\").replace('"', '\\"')
    return f'POSIX file "{escaped}"'


def _export_with_libreoffice(soffice_path: Path, input_path: Path, output_dir: Path, slide_count: int) -> None:
    if slide_count <= 0:
        return

    for page_number in range(1, slide_count + 1):
        page_output_dir = output_dir / f"page-{page_number}"
        page_output_dir.mkdir(parents=True, exist_ok=True)
        completed = subprocess.run(
            [
                str(soffice_path),
                "--headless",
                "--nologo",
                "--nodefault",
                "--nolockcheck",
                "--nofirststartwizard",
                "--convert-to",
                _build_libreoffice_png_filter(page_number),
                "--outdir",
                str(page_output_dir),
                str(input_path),
            ],
            capture_output=True,
            text=True,
            check=False,
        )
        if completed.returncode != 0:
            detail = (completed.stderr or completed.stdout or "").strip().splitlines()
            message = detail[-1] if detail else "未知错误"
            raise ApiError(code=500, msg=f"调用 LibreOffice 导出第 {page_number} 页失败：{message}", status_code=500)

        page_images = sorted(path for path in page_output_dir.iterdir() if path.is_file() and path.suffix.lower() == ".png")
        if len(page_images) != 1:
            raise ApiError(
                code=500,
                msg=f"LibreOffice 导出第 {page_number} 页异常，期望 1 张图片，实际 {len(page_images)} 张",
                status_code=500,
            )
        shutil.move(str(page_images[0]), str(output_dir / f"slide-{page_number}.png"))


def _build_libreoffice_png_filter(page_number: int) -> str:
    filter_options = {
        "PageNumber": {"type": "long", "value": str(page_number)},
    }
    return f"png:impress_png_Export:{json.dumps(filter_options, separators=(',', ':'))}"


def _render_pdf_to_png(pdf_path: Path, output_dir: Path, *, slide_count: int) -> None:
    try:
        with fitz.open(pdf_path) as document:
            page_total = document.page_count
            for page_index in range(page_total):
                page = document.load_page(page_index)
                pixmap = page.get_pixmap(alpha=False)
                target_path = output_dir / f"slide-{page_index + 1}.png"
                pixmap.save(target_path)
    except Exception as exc:  # noqa: BLE001
        raise ApiError(code=500, msg=f"PDF 渲染页图失败：{exc}", status_code=500) from exc

    exported_count = len([path for path in output_dir.iterdir() if path.is_file() and path.suffix.lower() == ".png"])
    if slide_count and exported_count != slide_count:
        raise ApiError(
            code=500,
            msg=f"PowerPoint for Mac 导出的 PDF 页数异常，期望 {slide_count} 页，实际 {exported_count} 页",
            status_code=500,
        )


def _clear_exported_images(export_dir: Path) -> None:
    for path in export_dir.iterdir():
        if path.is_file():
            path.unlink()
            continue
        if path.is_dir():
            shutil.rmtree(path)


def _collect_exported_images(export_dir: Path) -> list[Path]:
    png_files = [path for path in export_dir.rglob("*.png") if path.is_file()]
    if not png_files:
        return []

    numbered_files: list[tuple[int, Path]] = []
    unnumbered_files: list[Path] = []
    for path in png_files:
        # 兼容不同平台命名：Slide1 / 幻灯片 1 / page-1 等
        matches = re.findall(r"(\d+)", path.stem)
        if matches:
            numbered_files.append((int(matches[-1]), path))
        else:
            unnumbered_files.append(path)

    numbered_files.sort(key=lambda item: (item[0], str(item[1])))
    unnumbered_files.sort(key=str)
    return [path for _, path in numbered_files] + unnumbered_files


def _extract_title(shape: object) -> str | None:
    if not hasattr(shape, "is_placeholder") or not getattr(shape, "is_placeholder"):
        return None
    placeholder_format = getattr(shape, "placeholder_format", None)
    placeholder_type = getattr(placeholder_format, "type", None)
    placeholder_name = str(placeholder_type).lower() if placeholder_type is not None else ""
    if "title" not in placeholder_name and "center_title" not in placeholder_name and "subtitle" not in placeholder_name:
        return None
    text_frame = getattr(shape, "text_frame", None)
    if text_frame is None:
        return None
    text = _normalize_text(getattr(text_frame, "text", None))
    return text or None


def _extract_notes(slide: object) -> str | None:
    try:
        notes_slide = getattr(slide, "notes_slide")
    except Exception:  # noqa: BLE001
        return None

    text_frame = getattr(notes_slide, "notes_text_frame", None)
    if text_frame is None:
        return None
    text = _normalize_text(getattr(text_frame, "text", None))
    return text or None


def _normalize_text(value: str | None) -> str:
    if not value:
        return ""
    return "\n".join(line.strip() for line in value.splitlines() if line.strip())
