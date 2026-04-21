import tempfile
import unittest
from pathlib import Path
from unittest.mock import Mock, patch

from pptx import Presentation

from backend.app.common.exceptions import ApiError
from backend.app.parser.pptx_reader import extract_pptx_presentation


class PptxReaderTestCase(unittest.TestCase):
    temp_dir: tempfile.TemporaryDirectory[str] | None = None

    def setUp(self) -> None:
        self.temp_dir = tempfile.TemporaryDirectory()

    def tearDown(self) -> None:
        assert self.temp_dir is not None
        self.temp_dir.cleanup()

    @patch("backend.app.parser.pptx_reader._resolve_preview_exporters")
    def test_extract_pptx_presentation_returns_pages_and_previews(self, mock_resolve_preview_exporters) -> None:
        pptx_path = self._create_demo_pptx("demo.pptx")
        preview_dir = Path(self.temp_dir.name) / "previews"
        mock_resolve_preview_exporters.return_value = [("MockExporter", self._mock_exported_pngs)]

        file_info, extracted = extract_pptx_presentation(
            pptx_path.as_uri(),
            preview_output_dir=preview_dir,
            preview_public_base="/courseware-previews/parse-demo",
        )

        self.assertEqual(file_info.fileName, "demo.pptx")
        self.assertEqual(file_info.pageCount, 2)
        self.assertEqual(extracted.sourceType, "pptx")
        self.assertEqual(len(extracted.slides), 2)
        self.assertEqual(extracted.slides[0].slideNumber, 1)
        self.assertEqual(extracted.slides[0].previewUrl, "/courseware-previews/parse-demo/page-1.png")
        self.assertTrue((preview_dir / "page-1.png").exists())
        self.assertIn("第一章", extracted.slides[0].title or "")
        self.assertTrue(extracted.slides[0].bodyTexts)

    @patch("backend.app.parser.pptx_reader.platform.system", return_value="Windows")
    def test_resolve_preview_exporters_prefers_powerpoint_on_windows(self, _mock_system) -> None:
        from backend.app.parser import pptx_reader

        mock_powershell = Mock()
        mock_powershell.exists.return_value = True
        with patch.object(pptx_reader, "POWERSHELL_EXE", mock_powershell):
            exporters = pptx_reader._resolve_preview_exporters()

        self.assertEqual(exporters[0][0], "PowerPoint")
        self.assertIs(exporters[0][1], pptx_reader._export_with_powerpoint)

    @patch("backend.app.parser.pptx_reader._find_soffice_executable")
    @patch("backend.app.parser.pptx_reader.platform.system", return_value="Darwin")
    def test_resolve_preview_exporters_prefers_macos_powerpoint_before_libreoffice(self, _mock_system, mock_find_soffice) -> None:
        from backend.app.parser import pptx_reader

        mock_find_soffice.return_value = Path("/Applications/LibreOffice.app/Contents/MacOS/soffice")
        mock_powershell = Mock()
        mock_powershell.exists.return_value = False
        mock_osascript = Mock()
        mock_osascript.exists.return_value = True
        with patch.object(pptx_reader, "POWERSHELL_EXE", mock_powershell):
            with patch.object(pptx_reader, "OSASCRIPT_EXE", mock_osascript):
                with patch.object(
                    pptx_reader,
                    "_find_macos_powerpoint_app",
                    return_value=Path("/Applications/Microsoft PowerPoint.app"),
                ):
                    exporters = pptx_reader._resolve_preview_exporters()

        self.assertEqual(exporters[0][0], "PowerPoint for Mac")
        self.assertEqual(exporters[1][0], "LibreOffice")

    @patch("backend.app.parser.pptx_reader._find_soffice_executable", return_value=None)
    @patch("backend.app.parser.pptx_reader.platform.system", return_value="Darwin")
    def test_resolve_preview_exporters_raises_on_macos_without_powerpoint_or_libreoffice(self, _mock_system, _mock_find_soffice) -> None:
        from backend.app.parser import pptx_reader

        mock_powershell = Mock()
        mock_powershell.exists.return_value = False
        with patch.object(pptx_reader, "POWERSHELL_EXE", mock_powershell):
            with patch.object(pptx_reader, "_find_macos_powerpoint_app", return_value=None):
                with self.assertRaises(ApiError) as context:
                    pptx_reader._resolve_preview_exporters()

        self.assertIn("PowerPoint for Mac", context.exception.msg)

    def _mock_exported_pngs(self, _input_path: Path, export_dir: Path, _slide_count: int) -> None:
        (export_dir / "幻灯片1.PNG").write_bytes(b"slide-1")
        (export_dir / "幻灯片2.PNG").write_bytes(b"slide-2")

    def _create_demo_pptx(self, file_name: str) -> Path:
        assert self.temp_dir is not None
        target = Path(self.temp_dir.name) / file_name
        presentation = Presentation()

        title_slide = presentation.slides.add_slide(presentation.slide_layouts[0])
        title_slide.shapes.title.text = "第一章 绪论"
        title_slide.placeholders[1].text = "这是第一页的课程说明。"

        content_slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        content_slide.shapes.title.text = "第二页 内容"
        content_slide.placeholders[1].text = "这里是第二页正文。"

        presentation.save(target)
        return target
